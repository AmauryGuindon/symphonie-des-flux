# -*- coding: utf-8 -*-
"""
build_livrables.py
Génère les livrables DataCut :
  1. diagrams/*.png   — diagrammes vectoriels (matplotlib)
  2. Rapport_DataCut.docx
  3. Presentation_DataCut.pptx
La conversion PDF est faite ensuite via Word/PowerPoint COM (PowerShell).
"""
import os, json, collections

BASE = os.path.dirname(os.path.abspath(__file__))
DIAG = os.path.join(BASE, 'diagrams')
os.makedirs(DIAG, exist_ok=True)
DATASET = os.path.join(BASE, '..', 'DataCut', 'dataset create', 'dataset-v1-2026-04-13.json')

# ════════════════════════════════════════════════════════════════
# PALETTE
# ════════════════════════════════════════════════════════════════
C_DARK   = '#1A1A2E'
C_GOLD   = '#C9A44A'
C_BLUE   = '#105E8A'
C_GREEN  = '#2E7D32'
C_PURPLE = '#6A1B9A'
C_ORANGE = '#B75C00'
C_GREY   = '#555555'
C_LGREY  = '#888888'

# ════════════════════════════════════════════════════════════════
# 1. DIAGRAMMES (matplotlib)
# ════════════════════════════════════════════════════════════════
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

def _box(ax, cx, cy, w, h, text, fc, tc='white', fs=12, bold=True):
    x, y = cx - w/2, cy - h/2
    ax.add_patch(FancyBboxPatch((x, y), w, h,
        boxstyle="round,pad=0.6,rounding_size=2", linewidth=0,
        facecolor=fc, mutation_aspect=1))
    ax.text(cx, cy, text, ha='center', va='center', color=tc,
        fontsize=fs, fontweight='bold' if bold else 'normal', zorder=5,
        linespacing=1.35)

def _arrow(ax, x1, y1, x2, y2, color=C_GOLD, lw=2.2):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2),
        arrowstyle='-|>', mutation_scale=20, linewidth=lw,
        color=color, shrinkA=2, shrinkB=2, zorder=1))

def _label(ax, x, y, text, fs=10, color=C_GREY, ha='center', style='italic'):
    ax.text(x, y, text, ha=ha, va='center', fontsize=fs, color=color, style=style)

def _canvas(w_in, h_in):
    fig, ax = plt.subplots(figsize=(w_in, h_in))
    ax.set_xlim(0, 100); ax.set_ylim(0, 100); ax.axis('off')
    return fig, ax

def _save(fig, name):
    path = os.path.join(DIAG, name)
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white', pad_inches=0.15)
    plt.close(fig)
    print('  diagram:', name)
    return path

# ── architecture ────────────────────────────────────────────────
def diag_architecture():
    fig, ax = _canvas(11, 7)
    _box(ax, 50, 88, 42, 13, 'Frontend — Angular 17\nInterface admin + Interface client', C_BLUE, fs=13)
    _box(ax, 16, 88, 26, 13, 'MediaPipe\nDétection visage\n(in-browser, WASM)', C_ORANGE, fs=11)
    _label(ax, 16, 79, 'tourne dans\nle navigateur', fs=8, color=C_LGREY)
    # Frontend -> Backend
    _arrow(ax, 50, 81.5, 50, 70.5)
    _label(ax, 64, 76, 'HTTP REST  (4200 → 3001)', fs=10)
    _box(ax, 50, 64, 42, 13, 'Backend — NestJS\nAPI REST · Guards JWT · Multer', C_BLUE, fs=13)
    # storage
    _box(ax, 28, 40, 30, 12, 'MongoDB\nmétadonnées (datacut)', C_GREEN, fs=12)
    _box(ax, 72, 40, 30, 12, 'Disque local\n/uploads/gallery', C_PURPLE, fs=12)
    _arrow(ax, 44, 57.5, 32, 46.2)
    _arrow(ax, 56, 57.5, 68, 46.2)
    # python
    _box(ax, 50, 14, 42, 12, 'Python · scikit-learn\ntrain.py  →  model.pkl', C_DARK, fs=13)
    _arrow(ax, 28, 34, 44, 20.2)
    _label(ax, 28, 27, 'export\ndataset.json', fs=9)
    return _save(fig, 'architecture.png')

# ── pipeline 5 statuts ──────────────────────────────────────────
def diag_pipeline():
    fig, ax = _canvas(13, 3.4)
    stages = [
        ('RAW', 'Upload + features', C_GREY),
        ('VALIDATED', 'Vérif. qualité', C_BLUE),
        ('LABELED', 'Labels manuels', C_GOLD),
        ('PROCESSED', 'Validation finale', C_GREEN),
        ('EXPORTED', 'Intégré dataset', C_PURPLE),
    ]
    n = len(stages); w = 16; gap = (100 - n*w) / (n-1)
    for i, (name, desc, col) in enumerate(stages):
        cx = w/2 + i*(w+gap)
        _box(ax, cx, 60, w, 26, name, col, fs=13)
        _label(ax, cx, 36, desc, fs=10, color=C_GREY, style='normal')
        if i < n-1:
            _arrow(ax, cx + w/2, 60, cx + w/2 + gap, 60)
    _label(ax, 50, 12, 'Transitions unidirectionnelles — traçabilité complète à chaque étape',
           fs=11, color=C_GOLD)
    return _save(fig, 'pipeline.png')

# ── schéma BDD ──────────────────────────────────────────────────
def _entity(ax, x, title, fields, col):
    w = 28; top = 92; hh = 7
    ax.add_patch(FancyBboxPatch((x, top-hh), w, hh,
        boxstyle="round,pad=0.2,rounding_size=1", facecolor=col, linewidth=0))
    ax.text(x + w/2, top - hh/2, title, ha='center', va='center',
        color='white', fontsize=13, fontweight='bold')
    bh = 6.5 * len(fields) + 4
    ax.add_patch(FancyBboxPatch((x, top-hh-bh), w, bh,
        boxstyle="square,pad=0", facecolor='#F4F6F8',
        edgecolor=col, linewidth=1.5))
    for j, f in enumerate(fields):
        ax.text(x + 1.5, top - hh - 5 - j*6.5, f, ha='left', va='center',
            fontsize=10, color='#222222', family='monospace')
    return x + w

def diag_db():
    fig, ax = _canvas(13, 6.5)
    _entity(ax, 2, 'GalleryItem', [
        '_id        : ObjectId', 'filename   : String', 'url        : String',
        'status     : Enum', 'labels     : [String]', 'features   : {w,h,kb,color}',
        'datasetVersion : String?', 'createdAt  : Date',
    ], C_BLUE)
    _entity(ax, 36, 'DatasetVersion', [
        '_id            : ObjectId', 'version        : String', 'imageCount     : Number',
        'imageIds       : [String]', 'labelsIncluded : [String]', 'description    : String?',
        'createdAt      : Date',
    ], C_GREEN)
    _entity(ax, 70, 'User', [
        '_id          : ObjectId', 'email        : String', 'password     : bcrypt',
        'googleId     : String?', 'role         : Enum', 'loyaltyTier  : Enum',
        'createdAt    : Date',
    ], C_PURPLE)
    _arrow(ax, 30, 70, 36, 70, color=C_GOLD)
    _label(ax, 33, 76, 'datasetVersion\n→ version', fs=8, color=C_GOLD)
    _label(ax, 50, 6, 'Approche hybride : fichiers sur disque · métadonnées en MongoDB',
           fs=11, color=C_GREY)
    return _save(fig, 'db_schema.png')

# ── flux ML ─────────────────────────────────────────────────────
def diag_ml():
    fig, ax = _canvas(6.5, 9)
    steps = [
        ('dataset.json', C_GOLD),
        ('build_features() → vecteur 7D', '#222'),
        ('MultiLabelBinarizer → matrice y', '#222'),
        ('train_test_split (80 / 20)', '#222'),
        ('OneVsRest(RandomForest×100)', C_BLUE),
        ('classification_report + accuracy', '#222'),
        ('model.pkl  (pickle)', C_GOLD),
    ]
    n = len(steps); h = 9; gap = (100 - n*h) / (n-1)
    for i, (txt, col) in enumerate(steps):
        cy = 100 - h/2 - i*(h+gap)
        fc = col if col in (C_GOLD, C_BLUE) else '#E8EDF1'
        tc = 'white' if col in (C_GOLD, C_BLUE) else '#222222'
        _box(ax, 50, cy, 92, h, txt, fc, tc=tc, fs=12)
        if i < n-1:
            _arrow(ax, 50, cy - h/2, 50, cy - h/2 - gap, color=C_GREY)
    return _save(fig, 'ml_flow.png')

# ── flux recommandation ─────────────────────────────────────────
def diag_reco():
    fig, ax = _canvas(13, 3.4)
    steps = [
        ('Capture', 'Photo caméra\nou upload', C_BLUE),
        ('MediaPipe', '468 landmarks\ndétectés', C_ORANGE),
        ('Calcul', 'Forme du visage\n+ teinte de peau', C_GOLD),
        ('Résultat', 'Styles + galerie\nfiltrée', C_GREEN),
    ]
    n = len(steps); w = 19; gap = (100 - n*w) / (n-1)
    for i, (name, desc, col) in enumerate(steps):
        cx = w/2 + i*(w+gap)
        _box(ax, cx, 62, w, 30, name, col, fs=13)
        _label(ax, cx, 32, desc, fs=10, color=C_GREY, style='normal')
        if i < n-1:
            _arrow(ax, cx + w/2, 62, cx + w/2 + gap, 62)
    _label(ax, 50, 8, '100 % dans le navigateur — aucune donnée biométrique transmise (RGPD by design)',
           fs=11, color=C_GOLD)
    return _save(fig, 'reco_flow.png')

# ── distribution des labels (données réelles) ───────────────────
def diag_labels():
    with open(DATASET, encoding='utf-8') as f:
        data = json.load(f)
    counter = collections.Counter()
    for img in data['images']:
        counter.update(img['labels'])
    items = counter.most_common()
    labels = [k for k, _ in items][::-1]
    values = [v for _, v in items][::-1]
    fig, ax = plt.subplots(figsize=(8, 5))
    bars = ax.barh(labels, values, color=C_BLUE)
    bars[-1].set_color(C_GOLD)
    for i, v in enumerate(values):
        ax.text(v + 0.08, i, str(v), va='center', fontsize=10, color='#222', fontweight='bold')
    ax.set_xlabel("Nombre d'images", fontsize=11)
    ax.set_title(f"Distribution des labels — dataset v1 ({data['imageCount']} images)",
                 fontsize=13, fontweight='bold', color=C_DARK)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.set_xlim(0, max(values) + 1)
    fig.tight_layout()
    return _save(fig, 'labels_dist.png')

print('Génération des diagrammes…')
P_ARCH  = diag_architecture()
P_PIPE  = diag_pipeline()
P_DB    = diag_db()
P_ML    = diag_ml()
P_RECO  = diag_reco()
P_LABEL = diag_labels()

# ════════════════════════════════════════════════════════════════
# 2. RAPPORT WORD
# ════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

GOLD  = RGBColor(0xC9, 0xA4, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)
GREY  = RGBColor(0x55, 0x55, 0x55)

doc = Document()
for section in doc.sections:
    section.top_margin = Cm(2.5); section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.8); section.right_margin = Cm(2.3)

def set_font(run, size=11, bold=False, color=None, italic=False):
    run.font.name = 'Calibri'; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color

def h1(text):
    p = doc.add_heading(text, level=1)
    p.runs[0].font.color.rgb = DARK; p.runs[0].font.size = Pt(18); p.runs[0].font.bold = True

def h2(text):
    p = doc.add_heading(text, level=2)
    p.runs[0].font.color.rgb = RGBColor(0x10,0x5E,0x8A); p.runs[0].font.size = Pt(14); p.runs[0].font.bold = True

def h3(text):
    p = doc.add_heading(text, level=3)
    p.runs[0].font.color.rgb = RGBColor(0x20,0x7A,0xA8); p.runs[0].font.size = Pt(12); p.runs[0].font.bold = True

def body(text, bold=False, italic=False):
    p = doc.add_paragraph(); run = p.add_run(text)
    set_font(run, 11, bold=bold, italic=italic, color=RGBColor(0x22,0x22,0x22))
    p.paragraph_format.space_after = Pt(6); return p

def bullet(text):
    p = doc.add_paragraph(style='List Bullet'); run = p.add_run(text)
    set_font(run, 11); p.paragraph_format.space_after = Pt(3)

def code_block(text):
    p = doc.add_paragraph(); run = p.add_run(text)
    run.font.name = 'Consolas'; run.font.size = Pt(9); run.font.color.rgb = RGBColor(0x22,0x22,0x22)
    pf = p.paragraph_format; pf.left_indent = Cm(0.8); pf.space_before = Pt(4); pf.space_after = Pt(4)
    shd = OxmlElement('w:shd'); shd.set(qn('w:val'),'clear'); shd.set(qn('w:color'),'auto'); shd.set(qn('w:fill'),'F2F2F2')
    p._p.get_or_add_pPr().append(shd)

def figure(path, width=15.5, caption=None):
    doc.add_picture(path, width=Cm(width))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.paragraphs[-1].paragraph_format.space_before = Pt(6)
    if caption:
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(caption); set_font(r, 9, italic=True, color=GREY)
        p.paragraph_format.space_after = Pt(8)

def add_table(headers, rows):
    t = doc.add_table(rows=1+len(rows), cols=len(headers)); t.style = 'Table Grid'
    for i, hh in enumerate(headers):
        cell = t.rows[0].cells[i]; cell.text = hh
        for run in cell.paragraphs[0].runs:
            run.font.bold = True; run.font.color.rgb = WHITE; run.font.size = Pt(11)
        shd = OxmlElement('w:shd'); shd.set(qn('w:val'),'clear'); shd.set(qn('w:color'),'auto'); shd.set(qn('w:fill'),'105E8A')
        cell._tc.get_or_add_tcPr().append(shd)
    for ri, rd in enumerate(rows):
        fill = 'FFFFFF' if ri%2==0 else 'EEF4F8'
        for ci, val in enumerate(rd):
            cell = t.rows[ri+1].cells[ci]; cell.text = val
            for run in cell.paragraphs[0].runs: run.font.size = Pt(10)
            shd = OxmlElement('w:shd'); shd.set(qn('w:val'),'clear'); shd.set(qn('w:color'),'auto'); shd.set(qn('w:fill'),fill)
            cell._tc.get_or_add_tcPr().append(shd)
    doc.add_paragraph()

# ── PAGE DE GARDE ───────────────────────────────────────────────
doc.add_paragraph('\n\n\n')
p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r1 = p.add_run('DATA'); set_font(r1, 48, bold=True, color=DARK)
r2 = p.add_run('CUT');  set_font(r2, 48, bold=True, color=GOLD)
p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run('Symphonie des données en flux continu'); set_font(r, 16, italic=True, color=GREY)
doc.add_paragraph()
p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
for line, sz, bold in [
    ('Rapport technique', 13, True),
    ('\nMaster 1 Intelligence Artificielle — École 89', 12, False),
    ('\nIntervenant : Najib AL AWAR', 11, False),
    ('\nÉtudiant : Amaury Guindon', 11, False),
    ('\nAnnée : 2025–2026', 11, False),
]:
    r = p.add_run(line); set_font(r, sz, bold=bold, color=RGBColor(0x33,0x33,0x33))
doc.add_page_break()

# ── SOMMAIRE ────────────────────────────────────────────────────
h1('Sommaire')
for i, t in enumerate([
    'Introduction & contexte', 'Architecture de la plateforme', 'Schéma de base de données',
    'Pipeline de données', 'Extraction de features', 'Versioning des datasets',
    'Entraînement du modèle', 'Visualisation & supervision', 'Recommandation client par IA',
    'Bonnes pratiques & gouvernance', 'Conclusion & perspectives'], 1):
    bullet(f'{i}.  {t}')
doc.add_page_break()

# ── 1. INTRODUCTION ─────────────────────────────────────────────
h1('1. Introduction & contexte')
h2('1.1 Présentation du projet')
body("DataCut est une plateforme web de gestion de données conçue dans le cadre du projet « Symphonie des données en flux continu ». Elle s'inscrit dans le contexte fictif d'un barbershop haut de gamme souhaitant exploiter ses photos de réalisations pour entraîner un modèle de classification de styles de coupe.")
body("L'objectif central est de mettre en œuvre un cycle de vie complet des données : de l'ingestion d'images brutes jusqu'à l'alimentation continue d'un modèle d'intelligence artificielle, en passant par la labélisation, le versioning et l'export structuré.")
h2('1.2 Problématique')
body("Comment concevoir une infrastructure capable de collecter, structurer, labéliser et versionner des images, tout en garantissant la traçabilité complète du flux de données vers un modèle d'apprentissage automatique ?")
h2('1.3 Périmètre fonctionnel')
bullet("Backoffice (admin) : gestion du pipeline, labélisation, création et export de datasets versionnés, supervision via graphiques.")
bullet("Interface client : recommandation personnalisée de coupe par analyse IA du visage (forme et teinte de peau).")
doc.add_page_break()

# ── 2. ARCHITECTURE ─────────────────────────────────────────────
h1('2. Architecture de la plateforme')
h2('2.1 Vue d\'ensemble')
body("L'architecture repose sur une séparation claire en couches : un frontend Angular, un backend NestJS exposant une API REST, un stockage hybride (MongoDB pour les métadonnées, disque pour les fichiers) et un module Python d'entraînement.")
figure(P_ARCH, 15.5, 'Figure 1 — Architecture en couches de la plateforme DataCut')
h2('2.2 Stack technique')
add_table(['Couche','Technologie','Rôle'], [
    ['Frontend','Angular 17 (standalone, signals)','Interface utilisateur'],
    ['Backend','NestJS + Mongoose','API REST, logique métier'],
    ['Base de données','MongoDB','Stockage des métadonnées'],
    ['Stockage fichiers','Système de fichiers (multer)','Stockage des images brutes'],
    ['Feature extraction','sharp (Node.js)','Extraction automatique de caractéristiques'],
    ['ML entraînement','Python / scikit-learn','Classification multi-label'],
    ['ML client','MediaPipe (JavaScript)','Détection de forme de visage in-browser'],
])
h2('2.3 Justification des choix techniques')
bullet('MongoDB : flexibilité de schéma idéale pour des métadonnées hétérogènes, agrégations natives pour les statistiques de labels.')
bullet('NestJS : architecture modulaire (GalleryModule, DatasetModule, StatsModule), séparation claire des responsabilités.')
bullet("sharp : extraction de features légère côté serveur, sans dépendance externe, au moment de l'upload.")
bullet('scikit-learn : implémentation éprouvée du RandomForest multi-label, adaptée à un dataset de taille limitée.')
bullet('MediaPipe : s\'exécute en WebAssembly dans le navigateur, aucune donnée biométrique transmise à un serveur externe.')
doc.add_page_break()

# ── 3. BDD ──────────────────────────────────────────────────────
h1('3. Schéma de base de données')
h2('3.1 Principes de modélisation')
body("La base de données datacut suit une approche hybride : les fichiers images sont stockés sur disque (/uploads/gallery/), tandis que toutes les métadonnées associées (statut, labels, features, versioning) sont stockées dans MongoDB.")
figure(P_DB, 15.5, 'Figure 2 — Modèle de données : collections et relation logique')
h2('3.2 Collection GalleryItem')
code_block(
    "GalleryItem {\n"
    "  _id            : ObjectId      // Identifiant unique MongoDB\n"
    "  filename       : String        // Nom du fichier sur disque\n"
    "  url            : String        // Chemin (/uploads/gallery/...)\n"
    "  category       : String?       // coupe | barbe | dégradé...\n"
    "  status         : Enum          // raw|validated|labeled|processed|exported\n"
    "  labels         : [String]      // Labels assignés (fade, afro, tresse...)\n"
    "  features       : { width, height, sizeKb, format, dominantColor }\n"
    "  datasetVersion : String?       // Version d'export\n"
    "  createdAt, updatedAt : Date\n"
    "}")
h2('3.3 Collection DatasetVersion')
code_block(
    "DatasetVersion {\n"
    "  _id, version, imageCount, imageIds[], labelsIncluded[], description, createdAt\n"
    "}")
h2('3.4 Relations')
bullet("GalleryItem.datasetVersion → DatasetVersion.version (clé logique, document store).")
bullet("DatasetVersion.imageIds liste les GalleryItem inclus dans le snapshot.")
bullet("Relation volontairement légère : les snapshots sont immuables.")
doc.add_page_break()

# ── 4. PIPELINE ─────────────────────────────────────────────────
h1('4. Pipeline de données')
h2('4.1 Vue d\'ensemble')
body("Le pipeline structure le cycle de vie d'une image en 5 statuts successifs et unidirectionnels :")
figure(P_PIPE, 16.5, 'Figure 3 — Pipeline de données en 5 statuts')
h2('4.2 Description de chaque étape')
h3('Phase 1 — RAW (Ingestion)')
body("NestJS + Multer sauvegarde le fichier, sharp extrait automatiquement les features, et un document GalleryItem est créé avec status: 'raw'.")
h3('Phase 2 — VALIDATED')
body("L'administrateur vérifie la qualité visuelle (cadrage, netteté, pertinence) depuis l'interface Pipeline.")
h3('Phase 3 — LABELED')
body("Assignation de 1 à N labels parmi 12 disponibles (fade, dégradé, taper, burst-fade, afro, tresse, barbe, rasage, avant, après, enfant, adulte).")
h3('Phase 4 — PROCESSED')
body("Validation finale : l'image est complète et prête à intégrer un dataset.")
h3('Phase 5 — EXPORTED')
body("L'image est incluse dans un snapshot versionné ; son champ datasetVersion est renseigné.")
h2('4.3 Règles de transition')
bullet("Transitions unidirectionnelles : pas de retour à un statut antérieur.")
bullet("Seules les images labeled / processed / exported sont éligibles à un dataset.")
bullet("Le funnel Pipeline visualise les volumes par statut et révèle les goulots d'étranglement.")
doc.add_page_break()

# ── 5. FEATURES ─────────────────────────────────────────────────
h1('5. Extraction de features')
h2('5.1 Principe')
body("À chaque upload, le backend extrait automatiquement un vecteur de 7 caractéristiques numériques via sharp (Node.js). Ce vecteur constitue l'entrée du modèle ML, sans intervention humaine.")
h2('5.2 Features extraites')
add_table(['Index','Feature','Description','Exemple'], [
    ['0','width','Largeur en pixels','1080'],
    ['1','height','Hauteur en pixels','1350'],
    ['2','sizeKb','Taille du fichier en Ko','312.4'],
    ['3','ratio','Largeur / Hauteur','0.80'],
    ['4','r','Rouge de la couleur dominante','142'],
    ['5','g','Vert de la couleur dominante','98'],
    ['6','b','Bleu de la couleur dominante','67'],
])
h2('5.3 Implémentation Python')
code_block(
    "def build_features(images):\n"
    "    rows = []\n"
    "    for img in images:\n"
    "        feat  = img.get('features') or {}\n"
    "        w, h  = feat.get('width',0), feat.get('height',0)\n"
    "        size  = feat.get('sizeKb',0)\n"
    "        ratio = w / h if h > 0 else 1.0\n"
    "        r,g,b = hex_to_rgb(feat.get('dominantColor','#808080'))\n"
    "        rows.append([w, h, size, ratio, r, g, b])\n"
    "    return np.array(rows, dtype=float)")
doc.add_page_break()

# ── 6. VERSIONING ───────────────────────────────────────────────
h1('6. Versioning des datasets')
h2('6.1 Principe')
body("Le versioning repose sur des snapshots immuables : à un instant T, toutes les images éligibles sont capturées dans une version numérotée qui ne change jamais — garantissant la reproductibilité des entraînements.")
h2('6.2 Mécanisme de création')
bullet("L'administrateur clique sur « Créer une version ».")
bullet("Le backend identifie les images éligibles (labeled, processed, exported).")
bullet("Un document DatasetVersion est créé (numéro incrémental, imageIds, labelsIncluded, imageCount).")
bullet("Les images incluses passent à exported et reçoivent le tag de version.")
h2('6.3 Structure du JSON exporté')
code_block(
    '{\n'
    '  "version": "v1", "createdAt": "2026-04-13T10:20:16Z", "imageCount": 9,\n'
    '  "labelsIncluded": ["burst-fade","taper","barbe","afro","dégradé", ...],\n'
    '  "images": [\n'
    '    { "id":"...", "labels":["burst-fade","adulte","après"],\n'
    '      "features": {"width":705,"height":921,"sizeKb":677,"dominantColor":"#686878"} }\n'
    '  ]\n'
    '}')
doc.add_page_break()

# ── 7. ML ───────────────────────────────────────────────────────
h1('7. Entraînement du modèle')
h2('7.1 Algorithme choisi')
body("Le classificateur est un RandomForest multi-label (scikit-learn). La stratégie OneVsRestClassifier traite indépendamment chaque label : une image peut en porter plusieurs (ex : fade + barbe + adulte).")
h2('7.2 Pipeline d\'entraînement')
figure(P_ML, 9.5, 'Figure 4 — Chaîne d\'entraînement de train.py')
h2('7.3 Paramètres')
add_table(['Paramètre','Valeur','Justification'], [
    ['n_estimators','100','Compromis précision / vitesse'],
    ['test_size','0.2','Split 80/20 standard'],
    ['random_state','42','Reproductibilité'],
    ['Stratégie','OneVsRest','Indépendance des labels'],
])
h2('7.4 Utilisation du modèle')
code_block(
    "import pickle\n"
    "model = pickle.load(open('model.pkl','rb'))\n"
    "features = [[1080,1350,312.4,0.8,142,98,67]]\n"
    "pred   = model['clf'].predict(features)\n"
    "labels = model['mlb'].inverse_transform(pred)")
doc.add_page_break()

# ── 8. VISUALISATION ────────────────────────────────────────────
h1('8. Visualisation & supervision')
h2('8.1 Funnel Pipeline')
body("Un funnel interactif montre le volume d'images par statut. Chaque barre est cliquable et filtre la liste, révélant instantanément les goulots d'étranglement.")
h2('8.2 Distribution des labels')
body("La distribution des labels permet d'identifier les classes sous-représentées — indicateur clé de la qualité du dataset. Exemple sur le dataset v1 :")
figure(P_LABEL, 13, 'Figure 5 — Distribution réelle des labels du dataset v1')
h2('8.3 Croissance & statistiques de versions')
bullet("Graphique de croissance : nombre d'images ajoutées par semaine (8 semaines).")
bullet("Chaque version affiche : nombre d'images, labels inclus, date, export JSON direct.")
doc.add_page_break()

# ── 9. RECOMMANDATION ───────────────────────────────────────────
h1('9. Recommandation client par IA')
h2('9.1 Principe')
body("La page « Ma Coupe Idéale » fournit une recommandation personnalisée à partir du visage du client. L'analyse est réalisée entièrement dans le navigateur via MediaPipe — aucune donnée n'est transmise à un serveur externe.")
figure(P_RECO, 16.5, 'Figure 6 — Parcours de recommandation client (in-browser)')
h2('9.2 Détection de la forme du visage')
body("MediaPipe Face Landmarker détecte 468 points de repère. Les proportions calculées :")
add_table(['Mesure','Landmarks','Description'], [
    ['Hauteur du visage','10 → 152','Front → menton'],
    ['Largeur pommettes','234 → 454','Point le plus large'],
    ['Largeur mâchoire','172 → 397','Ligne mandibulaire'],
    ['Largeur front','70 → 300','Tempes'],
])
code_block(
    "ratio    = hauteur / largeur pommettes\n"
    "jawRatio = largeur mâchoire / largeur pommettes\n"
    "ratio > 1.60               → Oblongue\n"
    "ratio < 1.20 et jaw > 0.85 → Ronde\n"
    "front > 1.08 et jaw < 0.72 → Cœur\n"
    "jaw > 0.80 et ratio ≤ 1.55 → Carrée\n"
    "sinon                       → Ovale")
h2('9.3 Détection de la teinte de peau')
add_table(['Teinte','Luminosité HSL','Texture estimée','Styles exclus'], [
    ['Clair','> 58 %','Lisses à ondulés','afro, tresse'],
    ['Médium','36–58 %','Ondulés à bouclés','aucun'],
    ['Foncé','< 36 %','Crépus à bouclés serrés','aucun'],
])
doc.add_page_break()

# ── 10. GOUVERNANCE ─────────────────────────────────────────────
h1('10. Bonnes pratiques & gouvernance')
h2('10.1 Séparation fichiers / métadonnées')
bullet("Fichiers images → disque (/uploads/gallery/)")
bullet("Métadonnées → MongoDB (statut, labels, features, références)")
body("Cette séparation permet de migrer vers S3 ou un CDN sans toucher à la base de données.")
h2('10.2 Versioning immuable')
body("Les snapshots sont immuables. Relancer train.py sur le même JSON produit exactement le même modèle — reproductibilité et auditabilité garanties.")
h2('10.3 Sécurité & protection des données')
bullet("Authentification JWT, rôles client / admin, OAuth2 Google.")
bullet("Guards NestJS sur toutes les routes sensibles.")
bullet("MediaPipe in-browser : aucune donnée biométrique transmise (RGPD by design).")
doc.add_page_break()

# ── 11. CONCLUSION ──────────────────────────────────────────────
h1('11. Conclusion & perspectives')
h2('11.1 Bilan')
add_table(['Exigence','Statut'], [
    ['Architecture détaillée de la plateforme','Réalisé'],
    ['Schéma de base de données documenté','Réalisé'],
    ['Pipeline de données (5 statuts)','Réalisé'],
    ['Extraction de features automatique','Réalisé'],
    ['Versioning immuable des datasets','Réalisé'],
    ["Alimentation d'un modèle ML",'Réalisé'],
    ['Visualisation des flux et indicateurs','Réalisé'],
    ['Démonstration fonctionnelle','Réalisé'],
    ['Séparation fichiers / métadonnées','Réalisé'],
    ['Recommandation client IA (bonus)','Réalisé'],
])
h2('11.2 Perspectives d\'évolution')
bullet("Dataset plus volumineux (objectif : 100+ images par label).")
bullet("Stockage cloud (AWS S3, Cloudinary).")
bullet("API de prédiction temps réel exposant model.pkl.")
bullet("Réentraînement automatique déclenché par un nouveau dataset.")
bullet("CNN / Vision Transformer pour remplacer le RandomForest.")

OUT_DOCX = os.path.join(BASE, 'Rapport_DataCut.docx')
doc.save(OUT_DOCX)
print('Word généré :', OUT_DOCX)

# ════════════════════════════════════════════════════════════════
# 3. PRESENTATION POWERPOINT
# ════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.text import PP_ALIGN

BG    = PColor(0x0A,0x0A,0x0A)
GOLD2 = PColor(0xC9,0xA4,0x4A)
WHITE2= PColor(0xFF,0xFF,0xFF)
GREY2 = PColor(0xAA,0xAA,0xAA)
BLUE2 = PColor(0x10,0x5E,0x8A)

prs = Presentation(); prs.slide_width = PIn(13.33); prs.slide_height = PIn(7.5)
blank = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG; return s

def rect(s,l,t,w,h,color):
    sh = s.shapes.add_shape(1, PIn(l),PIn(t),PIn(w),PIn(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); return sh

def text(s,txt,l,t,w,h,size=18,bold=False,color=WHITE2,align=PP_ALIGN.LEFT,italic=False):
    tb = s.shapes.add_textbox(PIn(l),PIn(t),PIn(w),PIn(h)); tf = tb.text_frame; tf.word_wrap=True
    first=True
    for line in txt.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment = align; r = p.add_run(); r.text = line
        r.font.size=PPt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.italic=italic; r.font.name='Calibri'
    return tb

def titlebar(s,title,sub=None):
    rect(s,0,0,13.33,1.2,BLUE2)
    text(s,title,0.4,0.18,12.5,0.7,size=28,bold=True)
    if sub: text(s,sub,0.4,0.82,12.5,0.4,size=13,color=PColor(0xCC,0xDD,0xFF),italic=True)

def bullets(s,items,l,t,w,size=15,spacing=0.5,color=WHITE2):
    for i,it in enumerate(items):
        text(s,'▸  '+it,l,t+i*spacing,w,spacing+0.05,size=size,color=color)

def pic(s,path,l,t,w):
    s.shapes.add_picture(path,PIn(l),PIn(t),width=PIn(w))

# 1 — TITRE
s = slide()
rect(s,0,2.85,13.33,0.06,GOLD2)
text(s,'DATA',3.3,1.0,3.2,1.4,size=72,bold=True,align=PP_ALIGN.RIGHT)
text(s,'CUT',6.6,1.0,3.2,1.4,size=72,bold=True,color=GOLD2)
text(s,'Symphonie des données en flux continu',0,3.15,13.33,0.7,size=20,italic=True,color=GREY2,align=PP_ALIGN.CENTER)
text(s,'Rapport technique — Master 1 IA — École 89',0,3.95,13.33,0.5,size=14,color=GREY2,align=PP_ALIGN.CENTER)
text(s,'Amaury Guindon  ·  Intervenant : Najib AL AWAR  ·  2025–2026',0,4.55,13.33,0.5,size=12,color=PColor(0x77,0x77,0x77),align=PP_ALIGN.CENTER)

# 2 — CONTEXTE
s = slide(); titlebar(s,'Contexte & Problématique')
text(s,'Contexte',0.5,1.45,5.8,0.4,size=16,bold=True,color=GOLD2)
bullets(s,['Barbershop fictif exploitant ses photos de coupe',
           'Images labélisées : fade, afro, tresse, barbe…',
           'Besoin : pipeline structuré alimentant un modèle IA'],0.5,1.95,5.8)
text(s,'Problématique',7.0,1.45,5.8,0.4,size=16,bold=True,color=GOLD2)
text(s,"Comment collecter, labéliser et versionner des images tout en garantissant la traçabilité complète vers un modèle d'IA ?",7.0,1.95,5.8,1.5,size=14)
rect(s,0.5,5.5,12.3,0.06,GOLD2)
text(s,'Ingestion → Labélisation → Versioning → Entraînement → Recommandation',0.5,5.7,12.3,0.6,size=15,bold=True,color=GOLD2,align=PP_ALIGN.CENTER)

# 3 — ARCHITECTURE (image)
s = slide(); titlebar(s,'Architecture de la plateforme')
pic(s,P_ARCH,3.5,1.45,6.3)
text(s,'Séparation fichiers / métadonnées · ML offline + ML in-browser',0.5,6.7,12.3,0.5,size=13,color=GOLD2,align=PP_ALIGN.CENTER)

# 4 — STACK
s = slide(); titlebar(s,'Stack technique')
for i,(k,v) in enumerate([
    ('Frontend','Angular 17 — standalone, signals'),
    ('Backend','NestJS + Mongoose'),
    ('Base de données','MongoDB (base datacut)'),
    ('Feature extraction','sharp (Node.js)'),
    ('ML entraînement','Python · scikit-learn · RandomForest'),
    ('ML client','MediaPipe Face Landmarker (WebAssembly)'),
    ('Authentification','JWT + OAuth2 Google'),
]):
    y=1.45+i*0.74; rect(s,0.4,y,12.5,0.66,PColor(0x14,0x14,0x14) if i%2==0 else PColor(0x1E,0x1E,0x1E))
    text(s,k,0.6,y+0.1,3.6,0.5,size=14,bold=True,color=GOLD2)
    text(s,v,4.3,y+0.1,8.4,0.5,size=14)

# 5 — BDD (image)
s = slide(); titlebar(s,'Schéma de base de données')
pic(s,P_DB,1.4,1.45,10.5)
text(s,'GalleryItem.datasetVersion → DatasetVersion.version',0.5,6.9,12.3,0.4,size=12,italic=True,color=GOLD2,align=PP_ALIGN.CENTER)

# 6 — PIPELINE (image)
s = slide(); titlebar(s,'Pipeline de données — 5 statuts')
pic(s,P_PIPE,0.7,1.7,11.9)
bullets(s,['RAW — NestJS + Multer reçoit le fichier, sharp extrait les features',
           'VALIDATED — vérification visuelle de la qualité',
           'LABELED — 1 à N labels parmi 12 disponibles',
           'PROCESSED — validation finale avant export',
           'EXPORTED — image intégrée dans un snapshot versionné'],0.6,4.1,12.2,size=13,spacing=0.5)

# 7 — FEATURES
s = slide(); titlebar(s,'Extraction automatique de features (sharp)')
text(s,'À chaque upload — automatique, transparent pour l\'admin',0.5,1.35,12.3,0.45,size=13,color=GREY2,italic=True)
for i,(k,desc,ex) in enumerate([
    ('width','Largeur en pixels','1080'),('height','Hauteur en pixels','1350'),
    ('sizeKb','Taille du fichier','312 Ko'),('ratio','Largeur / Hauteur','0.80'),
    ('r, g, b','Couleur dominante (RGB)','142, 98, 67')]):
    y=1.95+i*0.74; rect(s,0.5,y,11.8,0.66,PColor(0x14,0x14,0x14) if i%2==0 else PColor(0x1C,0x1C,0x1C))
    text(s,k,0.7,y+0.1,2.5,0.5,size=14,bold=True,color=GOLD2)
    text(s,desc,3.3,y+0.1,6.0,0.5,size=14)
    text(s,ex,9.5,y+0.1,2.5,0.5,size=14,color=GREY2,align=PP_ALIGN.RIGHT)
text(s,'→ Vecteur 7D stocké en MongoDB · réutilisé directement par train.py',0.5,6.6,12.0,0.5,size=13,bold=True,color=GOLD2,align=PP_ALIGN.CENTER)

# 8 — VERSIONING
s = slide(); titlebar(s,'Versioning des datasets')
text(s,'Snapshots immuables · reproductibilité garantie',0.5,1.35,12.3,0.4,size=13,italic=True,color=GREY2)
for i,(num,desc) in enumerate([('1','Créer une version'),('2','Sélection images éligibles'),('3','Document DatasetVersion'),('4','Export JSON')]):
    x=0.8+i*3.1; rect(s,x,1.9,2.7,0.9,BLUE2)
    text(s,num,x+0.15,1.95,0.7,0.8,size=26,bold=True,color=GOLD2)
    text(s,desc,x+0.85,2.0,1.8,0.8,size=12)
bullets(s,['Versions numérotées (v1, v2…) et horodatées',
           'Les images conservent la référence à leur version',
           'Relancer train.py sur le même JSON → résultat identique',
           'Historique complet des versions dans l\'interface'],0.5,3.3,12.0,size=14,spacing=0.55)
rect(s,0.5,6.1,12.3,0.06,GOLD2)
text(s,'Contrainte du projet respectée : versioning immuable ✓',0.5,6.3,12.3,0.5,size=14,bold=True,color=GOLD2,align=PP_ALIGN.CENTER)

# 9 — ML (image)
s = slide(); titlebar(s,'Entraînement ML — train.py')
pic(s,P_ML,0.7,1.4,4.6)
text(s,'Stratégie multi-label',6.0,1.5,6.8,0.4,size=16,bold=True,color=GOLD2)
text(s,"OneVsRestClassifier permet à une image de porter plusieurs labels :\n\n   ex : fade + barbe + adulte",6.0,2.0,6.8,1.6,size=14)
text(s,'Paramètres',6.0,3.6,6.8,0.4,size=16,bold=True,color=GOLD2)
bullets(s,['n_estimators = 100','test_size = 0.2 (80/20)','random_state = 42'],6.0,4.1,6.8,size=13,spacing=0.45)
rect(s,6.0,5.7,6.8,0.6,PColor(0x14,0x14,0x14))
text(s,'py train.py dataset-v1.json',6.2,5.8,6.4,0.45,size=13,color=PColor(0x90,0xEE,0x90))

# 10 — RESULTATS (chart image)
s = slide(); titlebar(s,'Résultats & qualité du dataset')
pic(s,P_LABEL,0.6,1.5,6.6)
text(s,'Métriques générées',7.4,1.5,5.5,0.4,size=15,bold=True,color=GOLD2)
bullets(s,['Accuracy exacte','Precision / Recall / F1 par label','Support par label'],7.4,1.95,5.5,size=13,spacing=0.45)
text(s,'Contexte',7.4,3.5,5.5,0.4,size=15,bold=True,color=GOLD2)
bullets(s,['Dataset prototype : 9 images','Performances limitées par la taille','Architecture scalable sans toucher au code'],7.4,3.95,5.5,size=13,spacing=0.5)
rect(s,7.4,6.0,5.5,0.06,GOLD2)
text(s,'Plus de données → meilleur modèle',7.4,6.15,5.5,0.5,size=14,bold=True,color=GOLD2)

# 11 — RECOMMANDATION (image)
s = slide(); titlebar(s,'Recommandation client — « Ma Coupe Idéale »')
pic(s,P_RECO,0.7,1.45,11.9)
text(s,'Forme du visage',0.6,4.3,6.0,0.4,size=14,bold=True,color=GOLD2)
bullets(s,['Ovale → polyvalent','Ronde → fade, taper','Carrée → burst-fade, dégradé','Cœur → dégradé, taper','Oblongue → volume côtés'],0.6,4.75,5.9,size=12,spacing=0.36)
text(s,'Teinte de peau',7.0,4.3,5.8,0.4,size=14,bold=True,color=GOLD2)
bullets(s,['Clair → styles lisses (fade, taper)','Médium → tous les styles','Foncé → afro, tresse, burst-fade'],7.0,4.75,5.8,size=12,spacing=0.42)

# 12 — GOUVERNANCE
s = slide(); titlebar(s,'Bonnes pratiques & gouvernance')
for i,(t,d) in enumerate([
    ('✓ Séparation fichiers / métadonnées','Fichiers → disque · Métadonnées → MongoDB · migration S3/CDN possible'),
    ('✓ Versioning immuable','Snapshots non modifiables · reproductibilité · auditabilité'),
    ('✓ Sécurité & authentification','JWT + OAuth2 Google · rôles client/admin · Guards NestJS'),
    ('✓ Protection des données (RGPD)','MediaPipe in-browser · aucune donnée biométrique transmise')]):
    y=1.45+i*1.4; rect(s,0.4,y,12.5,1.2,PColor(0x12,0x12,0x12))
    text(s,t,0.7,y+0.12,12.0,0.45,size=15,bold=True,color=GOLD2)
    text(s,d,0.7,y+0.6,12.0,0.5,size=12,color=GREY2)

# 13 — CONCLUSION
s = slide(); titlebar(s,'Conclusion & Perspectives')
text(s,'Réalisé',0.5,1.45,5.8,0.4,size=16,bold=True,color=GOLD2)
bullets(s,['Pipeline complet (5 statuts)','Features automatiques (sharp)','Versioning immuable',
           'ML train.py → model.pkl','Visualisation & supervision','Recommandation IA (MediaPipe)',
           'Séparation fichiers / métadonnées'],0.5,1.95,5.8,size=13,spacing=0.52,color=WHITE2)
text(s,'Perspectives',7.0,1.45,5.8,0.4,size=16,bold=True,color=GOLD2)
bullets(s,['100+ images par label','Stockage cloud (S3, Cloudinary)','API REST exposant model.pkl',
           'Réentraînement automatique','CNN / Vision Transformer'],7.0,1.95,5.8,size=13,spacing=0.52,color=WHITE2)
rect(s,0.5,6.2,12.3,0.06,GOLD2)
text(s,'Architecture scalable : plus de données → meilleur modèle, sans modifier le code',0.5,6.4,12.3,0.6,size=14,bold=True,color=GOLD2,align=PP_ALIGN.CENTER)

OUT_PPTX = os.path.join(BASE, 'Presentation_DataCut.pptx')
prs.save(OUT_PPTX)
print('PowerPoint généré :', OUT_PPTX)
print('Terminé.')
